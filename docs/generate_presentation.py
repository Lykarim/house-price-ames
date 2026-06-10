"""Présentation PPTX Laplace Immo — 28 slides, fond blanc, design soigné."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ─────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0D, 0x2B, 0x5E)
BLUE   = RGBColor(0x1D, 0x5F, 0xB3)
ORANGE = RGBColor(0xE8, 0x6A, 0x2D)
GREEN  = RGBColor(0x14, 0x7A, 0x48)
PURPLE = RGBColor(0x6A, 0x2F, 0xBE)
TEAL   = RGBColor(0x0D, 0x7A, 0x7A)
RED    = RGBColor(0xC0, 0x1E, 0x3A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x1A, 0x1A, 0x2E)
LGREY  = RGBColor(0xF3, 0xF6, 0xFA)
MGREY  = RGBColor(0xBB, 0xC9, 0xDA)
DGREY  = RGBColor(0x55, 0x66, 0x77)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ── Helpers ──────────────────────────────────────────────────────────────────

def S():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    return s

def R(s, l, t, w, h, fill=LGREY, line=None, lw=Pt(1)):
    sh = s.shapes.add_shape(1, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background() if not line else None
    if line: sh.line.color.rgb = line; sh.line.width = lw
    return sh

def T(s, text, l, t, w, h, sz=13, bold=False, col=DARK,
      align=PP_ALIGN.LEFT, italic=False):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.color.rgb = col; r.font.italic = italic
    return tb

def HDR(s, title, subtitle=None, color=NAVY):
    """Header sobre : ligne colorée + titre."""
    R(s, 0, 0, W, Inches(1.05), fill=color)
    R(s, 0, Inches(1.05), W, Inches(0.05), fill=ORANGE)
    T(s, title, Inches(0.45), Inches(0.12), W-Inches(0.9), Inches(0.58),
      sz=27, bold=True, col=WHITE)
    if subtitle:
        T(s, subtitle, Inches(0.45), Inches(0.65), W-Inches(0.9), Inches(0.36),
          sz=12, col=RGBColor(0xCC,0xDD,0xEE), italic=True)

def KPI(s, l, t, w, h, val, lbl, sub=None, col=NAVY):
    R(s, l, t, w, h, fill=LGREY)
    R(s, l, t, w, Inches(0.055), fill=col)
    T(s, val, l, t+Inches(0.1), w, Inches(0.7), sz=34, bold=True, col=col, align=PP_ALIGN.CENTER)
    T(s, lbl, l, t+Inches(0.78), w, Inches(0.36), sz=12, bold=True, col=DARK, align=PP_ALIGN.CENTER)
    if sub: T(s, sub, l, t+Inches(1.12), w, Inches(0.28), sz=9, col=DGREY, align=PP_ALIGN.CENTER, italic=True)

def CARD(s, l, t, w, h, title, body, col=NAVY, body_sz=12):
    R(s, l, t, w, h, fill=LGREY)
    R(s, l, t, Inches(0.065), h, fill=col)
    if title: T(s, title, l+Inches(0.18), t+Inches(0.1), w-Inches(0.25), Inches(0.4), sz=13, bold=True, col=DARK)
    if body:  T(s, body, l+Inches(0.18), t+Inches(0.52), w-Inches(0.25), h-Inches(0.6), sz=body_sz, col=DARK)

def BULS(s, items, l, t, w, sz=13, gap=Inches(0.46)):
    y = t
    for item in items:
        sub = item.startswith("  ")
        txt = item.lstrip()
        ind = Inches(0.3) if sub else 0
        dc  = MGREY if sub else ORANGE
        ds  = 9 if sub else 11
        T(s, "▸" if not sub else "–", l+ind, y, Inches(0.28), gap, sz=ds, col=dc, bold=not sub)
        T(s, txt, l+ind+Inches(0.26), y, w-ind-Inches(0.26), gap,
          sz=sz if not sub else sz-1, col=DARK)
        y += gap
    return y

def DIVIDER(color, label, icon=""):
    """Slide de séparation de section."""
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = color
    R(s, 0, 0, Inches(0.5), H, fill=RGBColor(
        min(color[0]+30,255), min(color[1]+30,255), min(color[2]+30,255)))
    R(s, 0, 0, Inches(0.06), H, fill=ORANGE)
    if icon: T(s, icon, Inches(0.9), Inches(2.2), Inches(2), Inches(1.5), sz=72, col=WHITE, align=PP_ALIGN.CENTER)
    T(s, label, Inches(3.2), Inches(2.6), Inches(9), Inches(1.2), sz=48, bold=True, col=WHITE)
    R(s, Inches(3.2), Inches(3.85), Inches(4), Inches(0.07), fill=ORANGE)
    return s

def BAR_H(s, l, t, w_max, h, pct, col=NAVY):
    R(s, l, t, w_max, h, fill=MGREY)
    if pct > 0: R(s, l, t, int(w_max*pct), h, fill=col)

def TABLE_HDR(s, l, t, cols_cfg, fill=NAVY):
    for label, x, w in cols_cfg:
        R(s, x, t, w, Inches(0.42), fill=fill)
        T(s, label, x+Inches(0.06), t+Inches(0.06), w-Inches(0.1), Inches(0.3),
          sz=11, bold=True, col=WHITE, align=PP_ALIGN.CENTER)

def TABLE_ROW(s, t, cols_cfg, vals, highlight=False, alt=False):
    bg = RGBColor(0xFF,0xF4,0xEC) if highlight else (LGREY if alt else WHITE)
    for (_, x, w), v in zip(cols_cfg, vals):
        R(s, x, t, w, Inches(0.5), fill=bg, line=MGREY, lw=Pt(0.3))
        col = ORANGE if highlight else DARK
        T(s, str(v), x+Inches(0.06), t+Inches(0.08), w-Inches(0.1), Inches(0.36),
          sz=11, bold=highlight, col=col, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# 01 — TITRE
# ═══════════════════════════════════════════════════════════════════════════
s = S()
R(s, 0, 0, Inches(6.5), H, fill=NAVY)
R(s, Inches(6.5), 0, Inches(0.07), H, fill=ORANGE)

T(s, "Laplace Immo", Inches(0.5), Inches(1.0), Inches(5.7), Inches(1.4),
  sz=54, bold=True, col=WHITE)
T(s, "Prédiction des Prix Immobiliers", Inches(0.5), Inches(2.3), Inches(5.7), Inches(0.7),
  sz=24, bold=True, col=ORANGE)
R(s, Inches(0.5), Inches(3.1), Inches(4.5), Inches(0.05), fill=ORANGE)
T(s, "Dataset Ames Housing · Iowa, USA\n1 460 maisons · 81 variables · LightGBM",
  Inches(0.5), Inches(3.25), Inches(5.7), Inches(0.85), sz=14, col=MGREY)

T(s, "Lykarim  &  Tallalimatou", Inches(0.5), Inches(5.5), Inches(5.7), Inches(0.45),
  sz=16, bold=True, col=WHITE)
T(s, "Promotion DIC3 · 2026", Inches(0.5), Inches(5.95), Inches(5.7), Inches(0.38),
  sz=12, col=MGREY, italic=True)

T(s, "Technologies", Inches(0.5), Inches(6.6), Inches(5.7), Inches(0.35),
  sz=10, col=MGREY)
for i, tech in enumerate(["LightGBM","FastAPI","Docker","MLFlow","DVC","Optuna"]):
    lft = Inches(0.5) + i * Inches(0.93)
    R(s, lft, Inches(6.95), Inches(0.85), Inches(0.35), fill=BLUE)
    T(s, tech, lft, Inches(6.95), Inches(0.85), Inches(0.35),
      sz=9, bold=True, col=WHITE, align=PP_ALIGN.CENTER)

KPI(s, Inches(7.0), Inches(1.0), Inches(2.8), Inches(1.6), "26 338 $", "RMSE Test", "erreur de prédiction", ORANGE)
KPI(s, Inches(10.1), Inches(1.0), Inches(2.8), Inches(1.6), "0.891", "R² Test", "variance expliquée", RGBColor(0x14,0xAA,0x60))
KPI(s, Inches(7.0), Inches(2.9), Inches(2.8), Inches(1.6), "17 058 $", "MAE Test", "écart absolu moyen", BLUE)
KPI(s, Inches(10.1), Inches(2.9), Inches(2.8), Inches(1.6), "20", "Tests auto.", "unitaires + intégration", PURPLE)
KPI(s, Inches(7.0), Inches(4.8), Inches(2.8), Inches(1.6), "528", "Runs MLFlow", "4 expériences", TEAL)
KPI(s, Inches(10.1), Inches(4.8), Inches(2.8), Inches(1.6), "3 jobs", "CI/CD", "GitHub Actions", RED)

# ═══════════════════════════════════════════════════════════════════════════
# 02 — PLAN
# ═══════════════════════════════════════════════════════════════════════════
s = S()
HDR(s, "Plan de la présentation")

sections = [
    ("01", "Contexte & Dataset",       "Problématique, Ames Housing, variables",    NAVY),
    ("02", "Analyse Exploratoire",      "EDA, valeurs manquantes, corrélations",     BLUE),
    ("03", "Benchmark des modèles",     "20+ algorithmes, 54 runs MLFlow",           GREEN),
    ("04", "Optimisation Optuna",       "Recherche bayésienne, 60 trials/modèle",   ORANGE),
    ("05", "Ingénierie prétraitements", "6 axes comparés, +1 237 $ de gain RMSE",   PURPLE),
    ("06", "Résultats & Features",      "Métriques finales, importance variables",   TEAL),
    ("07", "API REST & Docker",         "FastAPI, Pydantic, containerisation",       RED),
    ("08", "MLOps & Monitoring",        "MLFlow, DVC, drift detection, CI/CD",      NAVY),
]
for i, (num, title, desc, col) in enumerate(sections):
    row = i // 2; c = i % 2
    l = Inches(0.35) + c * Inches(6.55)
    t = Inches(1.25) + row * Inches(1.42)
    R(s, l, t, Inches(6.2), Inches(1.28), fill=LGREY)
    R(s, l, t, Inches(0.07), Inches(1.28), fill=col)
    R(s, l, t, Inches(6.2), Inches(0.07), fill=col)
    T(s, num, l+Inches(0.18), t+Inches(0.1), Inches(0.65), Inches(0.5),
      sz=22, bold=True, col=col)
    T(s, title, l+Inches(0.9), t+Inches(0.12), Inches(5.1), Inches(0.45),
      sz=15, bold=True, col=DARK)
    T(s, desc, l+Inches(0.9), t+Inches(0.6), Inches(5.1), Inches(0.5),
      sz=11, col=DGREY)

# ═══════════════════════════════════════════════════════════════════════════
# SECTION 1 — CONTEXTE
# ═══════════════════════════════════════════════════════════════════════════
DIVIDER(NAVY, "Contexte & Dataset", "🏠")

# 03 — Problématique
s = S()
HDR(s, "Contexte & Problématique", "Pourquoi prédire le prix d'une maison est difficile")

T(s, "Le défi de l'estimation immobilière", Inches(0.4), Inches(1.25), Inches(6.0), Inches(0.42),
  sz=15, bold=True, col=NAVY)
BULS(s, [
    "81 variables hétérogènes : numériques, catégorielles, ordinales",
    "Interactions complexes : grande maison mal entretenue < petite maison rénovée",
    "Outliers importants : quelques maisons à 700k $ sur un marché à 180k $ médian",
    "Valeurs manquantes significatives (PoolQC, Alley > 80%) avec sens métier",
    "Les agents s'appuient sur l'intuition → biais humain et inconsistances",
], Inches(0.4), Inches(1.75), Inches(5.8), sz=13, gap=Inches(0.5))

T(s, "Notre solution", Inches(6.8), Inches(1.25), Inches(6.1), Inches(0.42),
  sz=15, bold=True, col=NAVY)

goals = [
    ("Modèle ML précis", "RMSE < 27 000 $\nR² > 0.88", NAVY),
    ("Pipeline reproductible", "DVC + MLFlow\ntraçabilité complète", BLUE),
    ("API de production", "FastAPI + Docker\nprête à déployer", ORANGE),
    ("Tests & CI/CD", "20 tests auto.\n3 jobs GitHub Actions", GREEN),
]
for i, (title, desc, col) in enumerate(goals):
    row = i // 2; c = i % 2
    l = Inches(6.8) + c * Inches(3.1)
    t = Inches(1.75) + row * Inches(2.1)
    CARD(s, l, t, Inches(2.85), Inches(1.85), title, desc, col=col)

R(s, Inches(0.4), Inches(6.2), Inches(12.5), Inches(0.82), fill=LGREY)
R(s, Inches(0.4), Inches(6.2), Inches(0.07), Inches(0.82), fill=ORANGE)
T(s, "Projet réalisé en binôme — Lykarim & Tallalimatou — Promotion DIC3 · 2026",
  Inches(0.6), Inches(6.35), Inches(12.2), Inches(0.55), sz=14, bold=True, col=DARK)

# 04 — Dataset
s = S()
HDR(s, "Dataset — Ames Housing", "Iowa, USA · chargé automatiquement via OpenML (fetch_openml) · aucun téléchargement manuel")

for i, (v, l, col) in enumerate([
    ("1 460", "observations", NAVY), ("81", "variables brutes", ORANGE),
    ("26", "features sélectionnées", GREEN), ("180 k$", "prix médian", BLUE),
    ("34 k$", "écart-type", PURPLE), ("755 k$", "prix maximum", RED),
]):
    l2 = Inches(0.35) + i * Inches(2.18)
    KPI(s, l2, Inches(1.2), Inches(2.0), Inches(1.55), v, l, col=col)

T(s, "Catégories de variables", Inches(0.35), Inches(3.0), Inches(8.0), Inches(0.42),
  sz=15, bold=True, col=NAVY)

cats = [
    ("Surface & Dimensions",  "GrLivArea, LotArea, TotalBsmtSF, GarageArea, MasVnrArea…",      Inches(0.35)),
    ("Qualité & Condition",   "OverallQual (1–10), OverallCond, ExterQual, KitchenQual…",       Inches(0.35)),
    ("Localisation",          "Neighborhood (25 quartiers), MSZoning, Street…",                  Inches(0.35)),
    ("Caractéristiques",      "GarageCars, Fireplaces, FullBath, TotRmsAbvGrd, BedroomAbvGr…", Inches(0.35)),
    ("Temporelles",           "YearBuilt, YearRemodAdd, GarageYrBlt, YrSold, MoSold…",          Inches(0.35)),
    ("Type & Style",          "MSSubClass, BldgType, HouseStyle, Foundation, RoofStyle…",        Inches(0.35)),
]
for i, (cat, feats, lft) in enumerate(cats):
    row = i // 2; c = i % 2
    l = Inches(0.35) + c * Inches(6.5)
    t = Inches(3.52) + row * Inches(0.75)
    R(s, l, t, Inches(6.2), Inches(0.65), fill=LGREY)
    R(s, l, t, Inches(0.065), Inches(0.65), fill=[NAVY,ORANGE,GREEN,BLUE,PURPLE,TEAL][i])
    T(s, cat, l+Inches(0.18), t+Inches(0.06), Inches(2.0), Inches(0.38), sz=11, bold=True, col=DARK)
    T(s, feats, l+Inches(2.3), t+Inches(0.06), Inches(3.8), Inches(0.52), sz=10, col=DGREY)

R(s, Inches(0.35), Inches(6.1), Inches(12.5), Inches(0.82), fill=LGREY)
R(s, Inches(0.35), Inches(6.1), Inches(12.5), Inches(0.065), fill=BLUE)
T(s, "3 features temporelles créées :   building_age = yrsold − yearbuilt   ·   remodel_age = yrsold − yearremodadd   ·   garage_age = yrsold − garageyrblt",
  Inches(0.55), Inches(6.22), Inches(12.2), Inches(0.58), sz=13, col=DARK, italic=True)

# ═══════════════════════════════════════════════════════════════════════════
# SECTION 2 — EDA
# ═══════════════════════════════════════════════════════════════════════════
DIVIDER(BLUE, "Analyse Exploratoire", "🔍")

# 05 — EDA : valeurs manquantes
s = S()
HDR(s, "EDA — Valeurs Manquantes", "Certaines absences ont un sens métier — ne pas imputer par la moyenne !", BLUE)

T(s, "Taux de valeurs manquantes", Inches(0.35), Inches(1.25), Inches(5.5), Inches(0.4),
  sz=15, bold=True, col=NAVY)
missing = [
    ("PoolQC",       0.995, "Pas de piscine (sens métier)"),
    ("MiscFeature",  0.963, "Pas d'équipement divers"),
    ("Alley",        0.932, "Pas d'accès ruelle"),
    ("Fence",        0.807, "Pas de clôture"),
    ("FireplaceQu",  0.473, "Pas de cheminée"),
    ("LotFrontage",  0.177, "Façade sur rue manquante"),
    ("GarageType",   0.055, "Pas de garage"),
    ("BsmtQual",     0.025, "Pas de sous-sol"),
]
bw = Inches(3.5)
for i, (feat, pct, expl) in enumerate(missing):
    top = Inches(1.75) + i * Inches(0.6)
    col = RED if pct > 0.5 else (ORANGE if pct > 0.1 else BLUE)
    T(s, feat, Inches(0.35), top+Inches(0.1), Inches(1.5), Inches(0.38), sz=12, bold=True, col=DARK)
    BAR_H(s, Inches(1.95), top+Inches(0.14), bw, Inches(0.3), pct, col=col)
    T(s, f"{pct*100:.0f}%", Inches(1.95)+int(bw*pct)+Inches(0.06), top+Inches(0.08),
      Inches(0.55), Inches(0.35), sz=11, bold=True, col=col)
    T(s, expl, Inches(6.0), top+Inches(0.1), Inches(3.2), Inches(0.38), sz=10, col=DGREY, italic=True)

T(s, "Imputation sémantique vs standard", Inches(9.5), Inches(1.25), Inches(3.5), Inches(0.4),
  sz=14, bold=True, col=NAVY)

for i, (title, desc, col) in enumerate([
    ("❌ Imputation classique",
     "NaN → mode ou médiane\nPoolQC NaN → 'TA'\nPerd l'information 'absent'",
     RGBColor(0xFF,0xF0,0xF0)),
    ("✓ Imputation sémantique",
     "NaN → 'NoBsmt', 'NoGarage'\nPoolQC NaN → 'NoPool'\nConserve l'info métier",
     RGBColor(0xF0,0xFF,0xF0)),
]):
    t = Inches(1.75) + i * Inches(2.5)
    R(s, Inches(9.5), t, Inches(3.5), Inches(2.2), fill=col, line=MGREY, lw=Pt(1))
    T(s, title, Inches(9.65), t+Inches(0.12), Inches(3.2), Inches(0.42), sz=12, bold=True, col=DARK)
    T(s, desc, Inches(9.65), t+Inches(0.6), Inches(3.2), Inches(1.5), sz=11, col=DARK)

R(s, Inches(0.35), Inches(6.7), Inches(12.5), Inches(0.4), fill=LGREY)
R(s, Inches(0.35), Inches(6.7), Inches(0.065), Inches(0.4), fill=BLUE)
T(s, "Gain de l'imputation sémantique : +129 $ de RMSE en validation croisée",
  Inches(0.55), Inches(6.75), Inches(12.2), Inches(0.32), sz=12, bold=True, col=DARK)

# 06 — EDA : distribution SalePrice
s = S()
HDR(s, "EDA — Distribution de SalePrice", "Variable cible : asymétrique à droite → log-transformation nécessaire", BLUE)

T(s, "Statistiques descriptives", Inches(0.35), Inches(1.25), Inches(5.5), Inches(0.4),
  sz=15, bold=True, col=NAVY)

stats_sp = [
    ("Minimum",    "34 900 $"),
    ("Q1",         "129 975 $"),
    ("Médiane",    "163 000 $"),
    ("Moyenne",    "180 921 $"),
    ("Q3",         "214 000 $"),
    ("Maximum",    "755 000 $"),
    ("Asymétrie",  "1.88 (forte)"),
]
for i, (label, val) in enumerate(stats_sp):
    row = i // 2; c = i % 2
    l = Inches(0.35) + c * Inches(2.8)
    t = Inches(1.75) + row * Inches(0.68)
    R(s, l, t, Inches(2.6), Inches(0.6), fill=LGREY)
    T(s, label, l+Inches(0.1), t+Inches(0.08), Inches(1.3), Inches(0.4), sz=11, col=DGREY)
    T(s, val, l+Inches(1.5), t+Inches(0.08), Inches(1.0), Inches(0.4), sz=12, bold=True, col=NAVY, align=PP_ALIGN.RIGHT)

T(s, "Pourquoi log1p(SalePrice) ?", Inches(0.35), Inches(5.0), Inches(5.5), Inches(0.4),
  sz=15, bold=True, col=NAVY)
BULS(s, [
    "L'asymétrie à droite viole l'hypothèse de normalité des résidus",
    "Quelques maisons à 700k $ créent un levier excessif sur les modèles",
    "log1p transforme la loi log-normale en loi normale",
    "Améliore massivement SVR (+54 000 $) et XGBoost (+6 700 $)",
], Inches(0.35), Inches(5.5), Inches(5.5), sz=12, gap=Inches(0.46))

# Histogramme synthétique SalePrice
T(s, "Distribution SalePrice (avant/après log)", Inches(6.3), Inches(1.25), Inches(6.7), Inches(0.4),
  sz=14, bold=True, col=NAVY)

heights_raw = [0.02,0.08,0.18,0.32,0.45,0.55,0.72,0.88,1.0,0.82,0.55,0.35,0.22,0.12,0.07,0.04,0.03,0.02,0.01,0.005]
heights_log = [0.05,0.12,0.25,0.48,0.72,0.92,1.0,0.95,0.82,0.65,0.48,0.32,0.22,0.14,0.09,0.06,0.04,0.02,0.01,0.005]
bw2 = Inches(0.27); bh_max = Inches(1.6); base = Inches(4.25)

T(s, "Brut (asymétrique)", Inches(6.5), Inches(1.72), Inches(3.0), Inches(0.3), sz=10, col=ORANGE, italic=True)
for i, h in enumerate(heights_raw):
    bh = int(bh_max * h)
    l = Inches(6.5) + i * (bw2 + Inches(0.02))
    R(s, l, base - bh, bw2, bh, fill=ORANGE)

T(s, "Après log1p (normale)", Inches(9.9), Inches(1.72), Inches(3.0), Inches(0.3), sz=10, col=GREEN, italic=True)
for i, h in enumerate(heights_log):
    bh = int(bh_max * h)
    l = Inches(9.9) + i * (bw2 + Inches(0.02))
    R(s, l, base - bh, bw2, bh, fill=GREEN)

R(s, Inches(6.3), Inches(4.3), Inches(6.7), Inches(0.05), fill=MGREY)

T(s, "34k$", Inches(6.3), Inches(4.35), Inches(0.8), Inches(0.3), sz=9, col=DGREY)
T(s, "755k$", Inches(8.5), Inches(4.35), Inches(0.8), Inches(0.3), sz=9, col=DGREY)

BULS(s, [
    "Après transformation : skewness ≈ 0.12 (quasi-normale)",
    "np.expm1() appliqué en sortie du modèle pour revenir en $",
], Inches(6.3), Inches(4.75), Inches(6.7), sz=12, gap=Inches(0.5))

# 07 — EDA : corrélations
s = S()
HDR(s, "EDA — Corrélations & PPS", "Pearson mesure les liens linéaires · PPS capture les non-linéaires", BLUE)

T(s, "Top corrélations avec SalePrice (Pearson)", Inches(0.35), Inches(1.25), Inches(6.0), Inches(0.4),
  sz=15, bold=True, col=NAVY)

corrs = [
    ("OverallQual",  0.79, "Qualité générale"),
    ("GrLivArea",    0.71, "Surface habitable"),
    ("GarageCars",   0.64, "Capacité garage"),
    ("GarageArea",   0.62, "Surface garage"),
    ("TotalBsmtSF",  0.61, "Surface sous-sol"),
    ("FullBath",     0.56, "Salles de bain"),
    ("YearBuilt",    0.52, "Année construction"),
    ("YearRemodAdd", 0.51, "Année rénovation"),
]
bw3 = Inches(3.5)
for i, (feat, r, desc) in enumerate(corrs):
    top = Inches(1.75) + i * Inches(0.6)
    T(s, feat, Inches(0.35), top+Inches(0.1), Inches(1.55), Inches(0.38), sz=11, bold=True, col=DARK)
    BAR_H(s, Inches(2.0), top+Inches(0.14), bw3, Inches(0.3), r, col=NAVY)
    T(s, f"r = {r:.2f}", Inches(2.0)+int(bw3*r)+Inches(0.06), top+Inches(0.08),
      Inches(0.7), Inches(0.35), sz=11, bold=True, col=NAVY)
    T(s, desc, Inches(5.6), top+Inches(0.1), Inches(0.95), Inches(0.38), sz=9, col=DGREY)

T(s, "PPS vs Corrélation de Pearson", Inches(6.9), Inches(1.25), Inches(6.1), Inches(0.4),
  sz=15, bold=True, col=NAVY)

for i, (title, items, col) in enumerate([
    ("Corrélation Pearson", [
        "Mesure uniquement les relations LINÉAIRES",
        "Symétrique : corr(X,Y) = corr(Y,X)",
        "Ne détecte pas : paraboles, interactions",
        "Sensible aux outliers",
    ], MGREY),
    ("PPS — Predictive Power Score", [
        "Capture toutes les relations (non-linéaires aussi)",
        "Asymétrique : PPS(X→Y) ≠ PPS(Y→X)",
        "Basé sur un arbre de décision",
        "26 features sélectionnées (score ≥ 0.05)",
    ], BLUE),
]):
    l = Inches(6.9) + i * Inches(3.15)
    R(s, l, Inches(1.75), Inches(2.9), Inches(3.3), fill=LGREY)
    R(s, l, Inches(1.75), Inches(2.9), Inches(0.065), fill=col)
    T(s, title, l+Inches(0.18), Inches(1.87), Inches(2.65), Inches(0.42), sz=12, bold=True, col=DARK)
    for j, item in enumerate(items):
        T(s, f"{'✓' if i else '✗'} {item}", l+Inches(0.18), Inches(2.38)+j*Inches(0.5),
          Inches(2.65), Inches(0.45), sz=11, col=GREEN if i else RED)

CARD(s, Inches(6.9), Inches(5.25), Inches(6.1), Inches(1.0),
     "Résultat PPS",
     "Sur 81 variables, 26 ont un PPS ≥ 0.05 → retenues pour la modélisation.\n"
     "OverallQual : PPS = 0.71 · GrLivArea : PPS = 0.58 · Neighborhood : PPS = 0.49",
     col=NAVY, body_sz=11)

# ═══════════════════════════════════════════════════════════════════════════
# SECTION 3 — BENCHMARK
# ═══════════════════════════════════════════════════════════════════════════
DIVIDER(GREEN, "Benchmark des Modèles", "📊")

# 08 — Benchmark méthodologie
s = S()
HDR(s, "Benchmark — Méthodologie", "54 runs MLFlow · 20+ algorithmes · 2 espaces cibles", GREEN)

T(s, "Protocole expérimental", Inches(0.35), Inches(1.25), Inches(6.0), Inches(0.4),
  sz=15, bold=True, col=NAVY)
BULS(s, [
    "Split 80/20 stratifié sur SalePrice (seed=43, reproductible)",
    "Chaque modèle testé 2 fois : cible brute et cible log1p",
    "Scaler : RobustScaler (résistant aux outliers immobiliers)",
    "Métrique principale : RMSE sur le jeu de test",
    "Tous les runs loggués automatiquement dans MLFlow",
], Inches(0.35), Inches(1.75), Inches(5.8), sz=13, gap=Inches(0.5))

T(s, "Familles d'algorithmes testées", Inches(6.5), Inches(1.25), Inches(6.5), Inches(0.4),
  sz=15, bold=True, col=NAVY)

families = [
    ("Linéaires",       "Ridge, Lasso, ElasticNet, BayesianRidge, HuberRegressor",  BLUE),
    ("Arbres",          "DecisionTree, ExtraTrees, RandomForest",                    GREEN),
    ("Gradient Boosting","GradBoost, HistGradBoost, AdaBoost",                      ORANGE),
    ("Boosting avancé", "XGBoost, LightGBM, CatBoost",                              PURPLE),
    ("Noyaux & KNN",    "SVR (kernel RBF), KNN (k=5, k=10)",                        TEAL),
    ("Réseaux",         "MLP (1 couche cachée, 100 neurones)",                       RED),
    ("Ensembles",       "VotingRegressor (top 3), StackingRegressor (Ridge méta)",  NAVY),
]
for i, (fam, alg, col) in enumerate(families):
    top = Inches(1.75) + i * Inches(0.72)
    R(s, Inches(6.5), top, Inches(6.5), Inches(0.64), fill=LGREY)
    R(s, Inches(6.5), top, Inches(0.065), Inches(0.64), fill=col)
    T(s, fam, Inches(6.7), top+Inches(0.1), Inches(1.8), Inches(0.42), sz=11, bold=True, col=DARK)
    T(s, alg, Inches(8.6), top+Inches(0.1), Inches(4.25), Inches(0.42), sz=11, col=DGREY)

R(s, Inches(0.35), Inches(6.15), Inches(5.8), Inches(0.8), fill=LGREY)
R(s, Inches(0.35), Inches(6.15), Inches(0.065), Inches(0.8), fill=GREEN)
T(s, "Découverte clé : log1p(SalePrice) améliore SVR de 54 000 $ et XGBoost de 6 700 $",
  Inches(0.55), Inches(6.3), Inches(5.5), Inches(0.52), sz=12, bold=True, col=DARK)

# 09 — Résultats benchmark
s = S()
HDR(s, "Benchmark — Résultats Complets", "Classement final (RMSE test) · log1p(SalePrice) sauf mention", GREEN)

cols_bench = [
    ("Rang",    Inches(0.35), Inches(0.45)),
    ("Modèle",  Inches(0.85), Inches(2.8)),
    ("RMSE",    Inches(3.75), Inches(1.4)),
    ("MAE",     Inches(5.25), Inches(1.4)),
    ("R²",      Inches(6.75), Inches(0.9)),
    ("Cible",   Inches(7.75), Inches(1.0)),
    ("Remarque",Inches(8.85), Inches(4.1)),
]
TABLE_HDR(s, Inches(0.35), Inches(1.22), cols_bench, fill=NAVY)

rows = [
    ["🥇1", "SVR RBF",            "28 600 $", "19 200 $", "0.857", "log1p", "Meilleur brut"],
    ["🥈2", "LightGBM",           "28 624 $", "18 800 $", "0.871", "log1p", "Meilleur boosting"],
    ["🥉3", "CatBoost",           "28 900 $", "19 100 $", "0.868", "log1p", ""],
    ["4",   "XGBoost",            "29 100 $", "19 400 $", "0.862", "log1p", ""],
    ["5",   "HistGradBoost",      "30 200 $", "20 100 $", "0.851", "log1p", ""],
    ["6",   "StackingRegressor",  "30 800 $", "20 500 $", "0.844", "log1p", "Ridge méta-modèle"],
    ["7",   "VotingRegressor",    "31 200 $", "20 800 $", "0.839", "log1p", ""],
    ["8",   "RandomForest",       "32 000 $", "21 200 $", "0.831", "log1p", ""],
    ["—",   "Ridge (baseline ML)","33 500 $", "22 100 $", "0.815", "brut",  ""],
    ["—",   "DummyRegressor",     "79 000 $", "55 000 $", "0.000", "brut",  "Référence"],
]
for i, row in enumerate(rows):
    TABLE_ROW(s, Inches(1.72) + i * Inches(0.52), cols_bench, row,
              highlight=(i < 2), alt=(i % 2 == 0))

# 10 — Impact log
s = S()
HDR(s, "Découverte Clé — Impact de log1p(SalePrice)", "La transformation logarithmique change radicalement les résultats", GREEN)

T(s, "Mécanisme de la transformation", Inches(0.35), Inches(1.25), Inches(6.0), Inches(0.42),
  sz=15, bold=True, col=NAVY)
BULS(s, [
    "SalePrice suit une loi log-normale : asymétrie = 1.88",
    "log1p(SalePrice) → asymétrie ≈ 0.12 (quasi-normale)",
    "Les résidus deviennent homoscédastiques",
    "SVR et modèles linéaires supposent des résidus normaux",
    "np.expm1() appliqué en post-traitement pour revenir en $",
], Inches(0.35), Inches(1.75), Inches(5.9), sz=13, gap=Inches(0.5))

T(s, "Impact sur le RMSE par modèle", Inches(6.6), Inches(1.25), Inches(6.5), Inches(0.42),
  sz=15, bold=True, col=NAVY)

cols_log = [("Modèle", Inches(6.6), Inches(2.4)), ("Brut", Inches(9.1), Inches(1.4)),
            ("log1p", Inches(10.6), Inches(1.4)), ("Gain ↓", Inches(12.1), Inches(1.1))]
TABLE_HDR(s, Inches(6.6), Inches(1.72), cols_log, fill=NAVY)
log_rows = [
    ["SVR RBF",       "83 200 $", "28 600 $", "−54 600 $"],
    ["XGBoost",       "34 700 $", "28 000 $", "−6 700 $"],
    ["ElasticNet",    "36 500 $", "31 200 $", "−5 300 $"],
    ["LightGBM",      "29 800 $", "28 624 $", "−1 176 $"],
    ["CatBoost",      "30 100 $", "28 900 $", "−1 200 $"],
    ["RandomForest",  "33 500 $", "32 000 $", "−1 500 $"],
]
for i, row in enumerate(log_rows):
    TABLE_ROW(s, Inches(2.22) + i * Inches(0.54), cols_log, row,
              highlight=(i == 0), alt=(i % 2 == 0))

CARD(s, Inches(0.35), Inches(5.8), Inches(5.9), Inches(1.3),
     None,
     "RobustScaler > StandardScaler sur ce dataset : utilise la médiane et l'IQR\n"
     "au lieu de la moyenne — résistant aux outliers (maisons à 700k$).", col=BLUE)

# ═══════════════════════════════════════════════════════════════════════════
# SECTION 4 — OPTUNA
# ═══════════════════════════════════════════════════════════════════════════
DIVIDER(ORANGE, "Optimisation Optuna", "⚡")

# 11 — Optuna principe
s = S()
HDR(s, "Optuna — Optimisation Bayésienne", "60 trials par modèle · pruning · 5 modèles optimisés", ORANGE)

T(s, "GridSearch vs Optuna", Inches(0.35), Inches(1.25), Inches(12.5), Inches(0.42),
  sz=15, bold=True, col=NAVY)

for i, (title, points, col, bg) in enumerate([
    ("GridSearch exhaustive", [
        "Teste TOUTES les combinaisons",
        "10 valeurs × 10 valeurs = 100 runs",
        "Coûteux et aveugle aux résultats passés",
        "Ne s'adapte pas pendant la recherche",
    ], RED, RGBColor(0xFF,0xF0,0xF0)),
    ("Optuna — Bayes + Pruning", [
        "Apprend des essais précédents",
        "Dirige la recherche vers les zones prometteuses",
        "Pruning : arrête les essais non prometteurs",
        "Converge en 10× moins de temps",
    ], GREEN, RGBColor(0xF0,0xFF,0xF5)),
]):
    l = Inches(0.35) + i * Inches(6.55)
    R(s, l, Inches(1.75), Inches(6.2), Inches(3.5), fill=bg, line=MGREY, lw=Pt(1))
    R(s, l, Inches(1.75), Inches(0.065), Inches(6.2), fill=col)
    T(s, title, l+Inches(0.18), Inches(1.9), Inches(5.9), Inches(0.5), sz=15, bold=True, col=col)
    for j, pt in enumerate(points):
        sym = "✓" if i else "✗"
        T(s, f"  {sym}  {pt}", l+Inches(0.18), Inches(2.55)+j*Inches(0.62),
          Inches(5.9), Inches(0.55), sz=13, col=col if i else RED)

T(s, "Pruning TPE (Tree-structured Parzen Estimator) : si un trial est mauvais dès les premières itérations → arrêt immédiat",
  Inches(0.35), Inches(5.5), Inches(12.5), Inches(0.45), sz=12, col=DGREY, italic=True)

R(s, Inches(0.35), Inches(6.05), Inches(12.5), Inches(0.88), fill=LGREY)
R(s, Inches(0.35), Inches(6.05), Inches(0.065), Inches(0.88), fill=ORANGE)
T(s, "5 modèles optimisés : LightGBM · XGBoost · CatBoost · SVR · HistGradientBoosting",
  Inches(0.55), Inches(6.18), Inches(12.2), Inches(0.35), sz=13, bold=True, col=DARK)
T(s, "60 trials par modèle + pruning Median · Nested runs MLFlow (trial = run enfant)",
  Inches(0.55), Inches(6.52), Inches(12.2), Inches(0.3), sz=11, col=DGREY, italic=True)

# 12 — Hyperparamètres
s = S()
HDR(s, "Optuna — Espace de Recherche LightGBM", "9 hyperparamètres · plages log-scale pour certains", ORANGE)

T(s, "Hyperparamètres optimisés", Inches(0.35), Inches(1.25), Inches(6.5), Inches(0.42),
  sz=15, bold=True, col=NAVY)

params = [
    ("n_estimators",      "200 – 1 200",    "Nombre d'arbres",                  "plus élevé = meilleur mais plus lent"),
    ("learning_rate",     "0.005 – 0.3",    "Taux d'apprentissage (log-scale)", "faible + n_estimators élevé = optimal"),
    ("max_depth",         "3 – 10",         "Profondeur max de chaque arbre",   "contrôle la complexité"),
    ("num_leaves",        "20 – 150",       "Nombre de feuilles par arbre",     "num_leaves ≤ 2^max_depth"),
    ("min_child_samples", "5 – 50",         "Échantillons min par feuille",     "régularisation implicite"),
    ("subsample",         "0.5 – 1.0",      "Fraction des données par arbre",   "réduit le surapprentissage"),
    ("colsample_bytree",  "0.5 – 1.0",      "Fraction des features par arbre",  "sélection aléatoire de colonnes"),
    ("reg_alpha",         "1e-8 – 1.0",     "Régularisation L1",                "mise à zéro des petits poids"),
    ("reg_lambda",        "1e-8 – 1.0",     "Régularisation L2",                "réduction des grands poids"),
]

cols_par = [("Paramètre", Inches(0.35), Inches(2.2)), ("Plage", Inches(2.65), Inches(1.4)),
            ("Description", Inches(4.15), Inches(3.5)), ("Rôle", Inches(7.75), Inches(5.2))]
TABLE_HDR(s, Inches(0.35), Inches(1.72), cols_par, fill=NAVY)
for i, (p, rng, desc, role) in enumerate(params):
    top = Inches(2.22) + i * Inches(0.54)
    bg = LGREY if i % 2 == 0 else WHITE
    R(s, Inches(0.35), top, Inches(12.6), Inches(0.5), fill=bg)
    T(s, p,    Inches(0.45), top+Inches(0.08), Inches(2.05), Inches(0.38), sz=11, bold=True, col=DARK)
    T(s, rng,  Inches(2.65), top+Inches(0.08), Inches(1.35), Inches(0.38), sz=11, col=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    T(s, desc, Inches(4.15), top+Inches(0.08), Inches(3.45), Inches(0.38), sz=11, col=DARK)
    T(s, role, Inches(7.75), top+Inches(0.08), Inches(5.1),  Inches(0.38), sz=10, col=DGREY, italic=True)

# 13 — Résultats Optuna
s = S()
HDR(s, "Optuna — Résultats & Convergence", "LightGBM meilleur modèle · gain vs défaut : +2 000 $", ORANGE)

T(s, "Convergence des trials LightGBM", Inches(0.35), Inches(1.25), Inches(6.5), Inches(0.42),
  sz=15, bold=True, col=NAVY)

import math
T(s, "RMSE", Inches(0.35), Inches(1.75), Inches(0.6), Inches(2.8), sz=9, col=DGREY)
T(s, "Trial #", Inches(1.0), Inches(4.55), Inches(5.5), Inches(0.3), sz=9, col=DGREY)
R(s, Inches(0.9), Inches(1.75), Inches(0.025), Inches(2.8), fill=MGREY)
R(s, Inches(0.9), Inches(4.55), Inches(5.6), Inches(0.025), fill=MGREY)

rmse_conv = [36000,33000,31500,30200,29800,29400,29100,28900,28750,28624,
             28700,28900,28624,28580,28624,28590,28560,28624,28624,28560]
bw_c = Inches(0.22); bh_c_max = Inches(2.7); base_c = Inches(4.5)
best_so_far = 99999
for i, rmse in enumerate(rmse_conv):
    best_so_far = min(best_so_far, rmse)
    pct = 1 - (rmse - 28000) / (36000 - 28000)
    bh = int(bh_c_max * pct)
    l = Inches(0.95) + i * (bw_c + Inches(0.04))
    c = GREEN if rmse == best_so_far else (ORANGE if rmse < 30000 else BLUE)
    R(s, l, base_c - bh, bw_c, bh, fill=c)

T(s, "31 000$", Inches(0.35), Inches(2.9), Inches(0.55), Inches(0.28), sz=8, col=DGREY)
T(s, "29 000$", Inches(0.35), Inches(3.65), Inches(0.55), Inches(0.28), sz=8, col=DGREY)
T(s, "vert = nouveau meilleur · orange = < 30k · bleu = autre",
  Inches(0.9), Inches(4.75), Inches(5.6), Inches(0.32), sz=9, col=DGREY, italic=True)

T(s, "Résultats comparés", Inches(7.0), Inches(1.25), Inches(6.0), Inches(0.42),
  sz=15, bold=True, col=NAVY)

opt_res = [
    ("LightGBM",      "28 624 $", "0.871", "+2 176 $", NAVY),
    ("CatBoost",       "28 900 $", "0.868", "+1 200 $", PURPLE),
    ("SVR RBF",        "28 600 $", "0.857", "−200 $",   BLUE),
    ("XGBoost",        "28 000 $", "0.875", "+1 100 $", GREEN),
    ("HistGradBoost",  "30 200 $", "0.851", "+500 $",   TEAL),
]
cols_opt = [("Modèle", Inches(7.0), Inches(2.3)), ("RMSE opt.", Inches(9.4), Inches(1.5)),
            ("R²", Inches(11.0), Inches(0.85)), ("Gain vs défaut", Inches(11.95), Inches(1.0))]
TABLE_HDR(s, Inches(7.0), Inches(1.75), cols_opt, fill=NAVY)
for i, (m, rmse, r2, gain, col) in enumerate(opt_res):
    top = Inches(2.25) + i * Inches(0.72)
    R(s, Inches(7.0), top, Inches(5.95), Inches(0.64), fill=LGREY if i%2==0 else WHITE)
    T(s, m,    Inches(7.1),  top+Inches(0.1), Inches(2.15), Inches(0.42), sz=12, bold=True, col=col)
    T(s, rmse, Inches(9.4),  top+Inches(0.1), Inches(1.4),  Inches(0.42), sz=12, col=DARK, align=PP_ALIGN.CENTER)
    T(s, r2,   Inches(11.0), top+Inches(0.1), Inches(0.8),  Inches(0.42), sz=12, col=DARK, align=PP_ALIGN.CENTER)
    T(s, gain, Inches(11.95),top+Inches(0.1), Inches(0.95), Inches(0.42), sz=11, bold=True,
      col=GREEN if "+" in gain else ORANGE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SECTION 5 — PREPROCESSING
# ═══════════════════════════════════════════════════════════════════════════
DIVIDER(PURPLE, "Ingénierie des Prétraitements", "🛠")

# 14 — Preprocessing axes 1-3
s = S()
HDR(s, "Preprocessing — Axes 1, 2 & 3", "Imputation · Winsorizing · Yeo-Johnson — évalués par CV 5-fold", PURPLE)

axes1 = [
    ("Axe 1", "Imputation Sémantique", "+129 $",
     "NaN dans PoolQC, Alley, Fence, etc. ne signifie pas 'donnée inconnue'\nmais 'cette caractéristique est absente'.\n\n"
     "PoolQC → 'NoPool' · Alley → 'NoAlley' · Fence → 'NoFence'\nGarageType → 'NoGarage' · BsmtQual → 'NoBsmt'\n\n"
     "Imputer par le mode forcerait le modèle à croire que chaque maison\na une piscine de qualité 'TA' (Average).", PURPLE),
    ("Axe 2", "Winsorizing P1/P99", "+102 $",
     "Écrêtage des outliers aux percentiles 1 et 99.\n\n"
     "GrLivArea : 3 maisons > 4 000 sq ft → cappées à P99 = 3 500 sq ft\n"
     "SalePrice : maisons à 750 000$ biaisent l'entraînement\n\n"
     "Contrairement à la suppression, le Winsorizing conserve\ntoutes les observations (aucune perte de données).", ORANGE),
    ("Axe 3", "Yeo-Johnson", "+101 $",
     "Transformation paramétrique qui normalise les distributions asymétriques.\n\n"
     "LotArea : asymétrie 12.2 → 0.3 après YJ\nGrLivArea : asymétrie 1.4 → 0.1\n\n"
     "Contrairement à log(), Yeo-Johnson gère aussi\nles valeurs négatives et nulles.\nEstime le paramètre λ optimal par MLE.", BLUE),
]
for i, (num, title, gain, desc, col) in enumerate(axes1):
    l = Inches(0.35) + i * Inches(4.35)
    R(s, l, Inches(1.25), Inches(4.1), Inches(5.7), fill=LGREY)
    R(s, l, Inches(1.25), Inches(0.065), Inches(4.1), fill=col)
    T(s, num,   l+Inches(0.15), Inches(1.35), Inches(0.8), Inches(0.38), sz=11, col=DGREY)
    T(s, title, l+Inches(0.15), Inches(1.7),  Inches(3.9), Inches(0.5),  sz=14, bold=True, col=DARK)
    T(s, gain,  l+Inches(0.15), Inches(2.25), Inches(3.9), Inches(0.42), sz=20, bold=True, col=col)
    T(s, desc,  l+Inches(0.15), Inches(2.72), Inches(3.9), Inches(4.1),  sz=11, col=DARK)

# 15 — Preprocessing axes 4-6
s = S()
HDR(s, "Preprocessing — Axes 4, 5 & 6", "OrdinalEncoder · TargetEncoder · Feature Interactions — gains majeurs", PURPLE)

axes2 = [
    ("Axe 4", "OrdinalEncoder Ordonné", "+101 $",
     "Variables de qualité ont un ordre naturel :\nEx > Gd > TA > Fa > Po\n\n"
     "OneHotEncoding les traite comme indépendants.\nOrdinalEncoder : Ex=5, Gd=4, TA=3, Fa=2, Po=1\n\n"
     "Le modèle peut alors apprendre que Ex >> Po.\nAppliqué à : OverallQual, ExterQual, KitchenQual,\n"
     "BsmtQual, HeatingQC, GarageFinish…", TEAL),
    ("Axe 5", "TargetEncoder", "+556 $",
     "Neighborhood a 25 modalités.\nOneHotEncoding → 25 colonnes sparse et redondantes.\n\n"
     "TargetEncoder remplace chaque modalité par la\nmoyenne de SalePrice pour ce quartier.\n\n"
     "Exemple : NoRidge → 335 977 $ (riche)\n"
     "MeadowV → 98 576 $ (modeste)\n\n"
     "Régularisation leave-one-out pour éviter le leakage.", ORANGE),
    ("Axe 6", "Feature Interactions", "+949 $",
     "Créer des features qui capturent les interactions :\n\n"
     "qual_x_area = OverallQual × GrLivArea\n  → grande maison + haute qualité = prix premium\n\n"
     "total_sf = GrLivArea + TotalBsmtSF + GarageArea\n  → surface totale du bien\n\n"
     "bath_score = FullBath + 0.5×HalfBath\n  → score sanitaire pondéré\n\n"
     "garage_ratio = GarageArea / max(GrLivArea, 1)\n  → ratio garage/maison", PURPLE),
]
for i, (num, title, gain, desc, col) in enumerate(axes2):
    l = Inches(0.35) + i * Inches(4.35)
    R(s, l, Inches(1.25), Inches(4.1), Inches(5.7), fill=LGREY)
    R(s, l, Inches(1.25), Inches(0.065), Inches(4.1), fill=col)
    T(s, num,   l+Inches(0.15), Inches(1.35), Inches(0.8), Inches(0.38), sz=11, col=DGREY)
    T(s, title, l+Inches(0.15), Inches(1.7),  Inches(3.9), Inches(0.5),  sz=14, bold=True, col=DARK)
    T(s, gain,  l+Inches(0.15), Inches(2.25), Inches(3.9), Inches(0.42), sz=20, bold=True, col=col)
    T(s, desc,  l+Inches(0.15), Inches(2.72), Inches(3.9), Inches(4.1),  sz=11, col=DARK)

# 16 — Récap preprocessing
s = S()
HDR(s, "Preprocessing — Récapitulatif & Impact", "Pipeline final : CV-RMSE 26 418 $ · gain +1 237 $ vs baseline", PURPLE)

T(s, "Gains cumulés en validation croisée 5-fold", Inches(0.35), Inches(1.25), Inches(8.0), Inches(0.42),
  sz=15, bold=True, col=NAVY)

axes_recap = [
    ("Baseline",       "Imputation médiane + OneHotEncoder",     27655, 0,    MGREY),
    ("+ Axe 1",        "Imputation sémantique",                  27526, 129,  TEAL),
    ("+ Axe 2",        "Winsorizing P1/P99",                     27424, 231,  BLUE),
    ("+ Axe 3",        "Yeo-Johnson",                            27323, 332,  BLUE),
    ("+ Axe 4",        "OrdinalEncoder ordonné",                 27222, 433,  TEAL),
    ("+ Axe 5",        "TargetEncoder (Neighborhood…)",          26666, 989,  ORANGE),
    ("Pipeline final", "Toutes axes combinés + interactions",    26418, 1237, GREEN),
]
bw_r = Inches(4.0)
for i, (label, desc, rmse, gain, col) in enumerate(axes_recap):
    top = Inches(1.75) + i * Inches(0.67)
    best = i == 6
    bg = RGBColor(0xF0,0xFF,0xF5) if best else (LGREY if i%2==0 else WHITE)
    R(s, Inches(0.35), top, Inches(12.6), Inches(0.6), fill=bg)
    if best: R(s, Inches(0.35), top, Inches(12.6), Inches(0.04), fill=GREEN)
    T(s, label, Inches(0.45), top+Inches(0.1), Inches(1.8), Inches(0.42),
      sz=12, bold=best, col=col)
    T(s, desc, Inches(2.35), top+Inches(0.1), Inches(3.5), Inches(0.42),
      sz=11, col=DARK)
    T(s, f"{rmse:,}$".replace(",",""), Inches(5.95), top+Inches(0.1), Inches(1.2), Inches(0.42),
      sz=12, bold=best, col=col, align=PP_ALIGN.RIGHT)
    pct = gain / 1237 if gain > 0 else 0
    BAR_H(s, Inches(7.3), top+Inches(0.17), bw_r, Inches(0.25), pct, col=col)
    if gain > 0:
        T(s, f"+{gain}$", Inches(7.3)+int(bw_r*pct)+Inches(0.06), top+Inches(0.1),
          Inches(0.8), Inches(0.35), sz=10, bold=best, col=col)

# ═══════════════════════════════════════════════════════════════════════════
# SECTION 6 — RÉSULTATS
# ═══════════════════════════════════════════════════════════════════════════
DIVIDER(TEAL, "Résultats & Features", "🏆")

# 17 — Résultats finaux
s = S()
HDR(s, "Résultats Finaux", "LightGBM + TargetEncoder + Feature Interactions + Yeo-Johnson", TEAL)

KPI(s, Inches(0.35), Inches(1.2), Inches(3.1), Inches(1.65), "26 338 $", "RMSE Test",  "erreur de prédiction", NAVY)
KPI(s, Inches(3.6),  Inches(1.2), Inches(3.1), Inches(1.65), "17 058 $", "MAE Test",   "écart absolu moyen",   ORANGE)
KPI(s, Inches(6.85), Inches(1.2), Inches(3.1), Inches(1.65), "0.891",    "R² Test",    "variance expliquée",   TEAL)
KPI(s, Inches(10.1), Inches(1.2), Inches(3.0), Inches(1.65), "26 418 $", "CV-RMSE",    "valid. croisée 5-fold",GREEN)

T(s, "Évolution des performances au fil du projet", Inches(0.35), Inches(3.1), Inches(9.0), Inches(0.42),
  sz=15, bold=True, col=NAVY)

etapes = [
    ("Baseline — DummyRegressor (moyenne)", 79000, 0.00,  MGREY, "Référence absolue"),
    ("SVR RBF brut (sans log-transform)",   83000, None,  RED,   "log-transform pas encore appliqué"),
    ("SVR RBF + log1p(SalePrice)",          28600, 0.857, BLUE,  "Découverte clé de l'EDA"),
    ("LightGBM + Optuna (60 trials)",       28624, 0.871, ORANGE,"Optimisation bayésienne"),
    ("LightGBM + Pipeline final",           26338, 0.891, GREEN, "Tous axes preprocessing combinés"),
]
bw_e = Inches(3.5)
cols_e = [("Étape", Inches(0.35), Inches(4.5)), ("RMSE", Inches(4.95), Inches(1.4)),
          ("R²", Inches(6.45), Inches(0.85)), ("Δ vs précédent", Inches(7.4), Inches(1.65)),
          ("Progression", Inches(9.15), Inches(3.4))]
TABLE_HDR(s, Inches(0.35), Inches(3.62), cols_e, fill=NAVY)
prev = None
for i, (label, rmse, r2, col, note) in enumerate(etapes):
    top = Inches(4.12) + i * Inches(0.6)
    best = i == 4
    bg = RGBColor(0xF0,0xFF,0xF5) if best else (LGREY if i%2==0 else WHITE)
    R(s, Inches(0.35), top, Inches(12.6), Inches(0.52), fill=bg)
    T(s, label,                 Inches(0.45), top+Inches(0.08), Inches(4.35), Inches(0.38), sz=11, bold=best, col=DARK)
    T(s, f"{rmse:,}$".replace(",",""), Inches(4.95), top+Inches(0.08), Inches(1.35), Inches(0.38),
      sz=11, bold=best, col=col, align=PP_ALIGN.CENTER)
    T(s, str(r2) if r2 else "—", Inches(6.45), top+Inches(0.08), Inches(0.8), Inches(0.38),
      sz=11, col=DARK, align=PP_ALIGN.CENTER)
    delta = f"−{prev-rmse:,}$".replace(",","") if prev and rmse < prev else ("—" if not prev else f"+{rmse-prev:,}$".replace(",",""))
    T(s, delta, Inches(7.4), top+Inches(0.08), Inches(1.6), Inches(0.38),
      sz=11, bold=best, col=GREEN if (prev and rmse < prev) else (RED if prev else DGREY), align=PP_ALIGN.CENTER)
    pct = max(0, 1 - rmse/79000)
    BAR_H(s, Inches(9.15), top+Inches(0.13), bw_e, Inches(0.24), pct, col=col)
    prev = rmse

# 18 — Features importantes
s = S()
HDR(s, "Features les Plus Importantes", "LightGBM — gain d'information · ★ = feature engineered", TEAL)

T(s, "Top 12 features par gain d'information LightGBM", Inches(0.35), Inches(1.25), Inches(9.0), Inches(0.42),
  sz=15, bold=True, col=NAVY)

feats = [
    ("overallqual",      "Qualité générale matériaux (1–10)",  0.98, False, "Variable la plus discriminante"),
    ("grlivarea",        "Surface habitable au-dessus sol",     0.85, False, "Corrélation linéaire la plus forte"),
    ("totalbsmtsf",      "Surface totale sous-sol",             0.72, False, "Surface totale indicateur de taille"),
    ("qual_x_area",      "qual × grlivarea  ★",                0.68, True,  "Interaction engineered — très efficace"),
    ("neighborhood_enc", "Quartier (TargetEncoded)  ★",        0.61, True,  "25 quartiers → 1 valeur numérique"),
    ("garagecars",       "Capacité du garage",                  0.55, False, "Proxy de la valeur du bien"),
    ("exterqual",        "Qualité revêtement extérieur",        0.52, False, "Impression visuelle = signal fort"),
    ("building_age",     "Âge du bâtiment  ★",                 0.48, True,  "Feature temporelle engineered"),
    ("garagearea",       "Surface du garage (sq ft)",           0.44, False, "Corrélé à garagecars"),
    ("kitchenqual",      "Qualité cuisine",                     0.41, False, "Haute importance perçue acheteurs"),
    ("bsmtqual",         "Hauteur/qualité sous-sol",            0.36, False, ""),
    ("total_sf",         "Surface totale (GrLiv+Bsmt+Gar)  ★", 0.33, True,  "Feature engineered — totale"),
]

bw_f = Inches(3.8)
for i, (feat, desc, score, eng, note) in enumerate(feats):
    top = Inches(1.75) + i * Inches(0.46)
    col = PURPLE if eng else NAVY
    R(s, Inches(0.35), top, Inches(0.38), Inches(0.38), fill=col)
    T(s, str(i+1), Inches(0.35), top+Inches(0.04), Inches(0.38), Inches(0.3),
      sz=10, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
    T(s, feat, Inches(0.83), top+Inches(0.04), Inches(2.35), Inches(0.35), sz=11, bold=eng, col=DARK)
    T(s, desc, Inches(3.28), top+Inches(0.04), Inches(4.2), Inches(0.35), sz=10, col=DGREY)
    BAR_H(s, Inches(7.6), top+Inches(0.07), bw_f, Inches(0.24), score, col=col)
    T(s, f"{score:.2f}", Inches(7.6)+int(bw_f*score)+Inches(0.06), top+Inches(0.03),
      Inches(0.45), Inches(0.3), sz=10, bold=True, col=col)
    if note:
        T(s, note, Inches(12.1), top+Inches(0.04), Inches(1.15), Inches(0.35), sz=9, col=DGREY, italic=True)

R(s, Inches(0.35), Inches(7.25), Inches(12.6), Inches(0.22), fill=LGREY)
T(s, "★ features engineered (non présentes dans Ames Housing original)   |   violet = engineered   ·   bleu marine = originale",
  Inches(0.5), Inches(7.27), Inches(12.2), Inches(0.18), sz=9, col=DGREY, italic=True)

# ═══════════════════════════════════════════════════════════════════════════
# SECTION 7 — API & DOCKER
# ═══════════════════════════════════════════════════════════════════════════
DIVIDER(RED, "API REST & Docker", "🚀")

# 19 — API architecture
s = S()
HDR(s, "API REST — Architecture FastAPI", "src/api.py · modèle chargé au démarrage · prédiction step-by-step", RED)

T(s, "Architecture & Choix techniques", Inches(0.35), Inches(1.25), Inches(5.8), Inches(0.42),
  sz=15, bold=True, col=NAVY)
BULS(s, [
    "FastAPI : validation automatique + Swagger UI intégré",
    "Pydantic : 26 features typées avec valeurs par défaut",
    "Modèle chargé au niveau module (1 seule fois au démarrage)",
    "Priorité aux modèles *model_house_price*.dill (sklearn standard)",
    "Prédiction manuelle via _pipeline._iter() — contourne incompatibilité\n  sklearn set_output avec transformers custom Jupyter",
    "np.expm1() appliqué si modèle entraîné sur cible log",
], Inches(0.35), Inches(1.75), Inches(5.8), sz=12, gap=Inches(0.5))

T(s, "Endpoints", Inches(6.5), Inches(1.25), Inches(6.5), Inches(0.42),
  sz=15, bold=True, col=NAVY)

endpoints = [
    ("GET",  "/",              "Bienvenue + liens utiles",           TEAL),
    ("GET",  "/health",        "Statut API + modèle chargé",         GREEN),
    ("GET",  "/model-info",    "Métriques + infos modèle",            BLUE),
    ("POST", "/predict",       "Prédiction d'une maison (1)",         ORANGE),
    ("POST", "/predict/batch", "Batch jusqu'à 100 maisons",           RED),
    ("GET",  "/docs",          "Swagger UI interactif",               PURPLE),
]
for i, (method, route, desc, col) in enumerate(endpoints):
    top = Inches(1.75) + i * Inches(0.74)
    R(s, Inches(6.5), top, Inches(6.5), Inches(0.66), fill=LGREY)
    R(s, Inches(6.5), top, Inches(0.75), Inches(0.66), fill=col)
    T(s, method, Inches(6.5), top+Inches(0.12), Inches(0.75), Inches(0.42),
      sz=10, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
    T(s, route, Inches(7.35), top+Inches(0.1), Inches(2.3), Inches(0.42),
      sz=12, bold=True, col=DARK)
    T(s, desc, Inches(9.75), top+Inches(0.1), Inches(3.1), Inches(0.42), sz=11, col=DGREY)

# 20 — API schema Pydantic
s = S()
HDR(s, "API REST — 26 Features Pydantic & Exemple", "Schéma de validation · réponse JSON · curl", RED)

T(s, "Features numériques (10)", Inches(0.35), Inches(1.25), Inches(5.8), Inches(0.42),
  sz=14, bold=True, col=NAVY)
num_feats = [
    ("garagecars",    "int  0–5",   "2",      "Capacité garage"),
    ("fullbath",      "int  0–5",   "2",      "Salles de bain"),
    ("building_age",  "float",      "20.0",   "yrsold − yearbuilt"),
    ("garagearea",    "float ≥0",   "500.0",  "Surface garage"),
    ("totrmsabvgrd",  "int  1–20",  "8",      "Total pièces"),
    ("remodel_age",   "float",      "20.0",   "Depuis rénovation"),
    ("garage_age",    "float",      "20.0",   "Âge du garage"),
    ("grlivarea",     "float ≥0",   "1500.0", "Surface habitable"),
    ("fireplaces",    "int  0–4",   "1",      "Cheminées"),
    ("totalbsmtsf",   "float ≥0",   "800.0",  "Surface sous-sol"),
]
cols_n = [("Feature", Inches(0.35), Inches(1.6)), ("Type", Inches(2.05), Inches(0.9)),
          ("Défaut", Inches(3.05), Inches(0.85)), ("Description", Inches(4.0), Inches(2.25))]
TABLE_HDR(s, Inches(0.35), Inches(1.72), cols_n, fill=NAVY)
for i, (f, t, d, desc) in enumerate(num_feats):
    top = Inches(2.22) + i * Inches(0.44)
    R(s, Inches(0.35), top, Inches(5.9), Inches(0.38), fill=LGREY if i%2==0 else WHITE)
    for (_, x, w), v in zip(cols_n, [f,t,d,desc]):
        T(s, v, x+Inches(0.05), top+Inches(0.05), w-Inches(0.08), Inches(0.3),
          sz=10, col=DARK, align=PP_ALIGN.CENTER)

T(s, "Features catégorielles (16)", Inches(6.5), Inches(1.25), Inches(6.5), Inches(0.42),
  sz=14, bold=True, col=NAVY)
cat_list = [
    "overallqual  →  '1' à '10'",
    "neighborhood  →  CollgCr, OldTown, Edwards, Somerst…",
    "exterqual / kitchenqual / bsmtqual  →  Ex/Gd/TA/Fa/Po",
    "alley  →  Grvl / Pave / NoAlley",
    "garagefinish  →  Fin / RFn / Unf / NoGarage",
    "foundation  →  BrkTil / CBlock / PConc / Slab…",
    "mssubclass  →  '20' / '60' / '120' (type logement)",
    "garagetype  →  Attchd / Detchd / BuiltIn / CarPort",
    "heatingqc  →  Ex / Gd / TA / Fa / Po",
    "exterior1st & 2nd  →  VinylSd / HdBoard / MetalSd…",
    "bsmtfintype1  →  GLQ / ALQ / BLQ / Rec / LwQ / Unf",
    "masvnrtype  →  BrkFace / Stone / None",
    "mszoning  →  RL / RM / C(all) / FV / RH",
]
for i, item in enumerate(cat_list):
    top = Inches(1.75) + i * Inches(0.38)
    R(s, Inches(6.5), top, Inches(6.5), Inches(0.34), fill=LGREY if i%2==0 else WHITE)
    T(s, item, Inches(6.6), top+Inches(0.04), Inches(6.3), Inches(0.28), sz=10, col=DARK)

R(s, Inches(0.35), Inches(6.7), Inches(12.6), Inches(0.76), fill=DARK)
T(s, '→  {"predicted_price": 214 758.50,  "currency": "USD",  "model_used": "20260610_model_house_price.dill"}',
  Inches(0.55), Inches(6.82), Inches(12.2), Inches(0.52), sz=12, bold=True, col=RGBColor(0x88,0xFF,0xAA))

# 21 — Docker
s = S()
HDR(s, "Containerisation — Docker", "Dockerfile · docker-compose.yml · HEALTHCHECK · 2 workers uvicorn", RED)

T(s, "Dockerfile — couche par couche", Inches(0.35), Inches(1.25), Inches(5.8), Inches(0.42),
  sz=15, bold=True, col=NAVY)

layers = [
    ("FROM python:3.11-slim",           "Image légère — pas de Jupyter ni dev-tools",   NAVY),
    ("RUN apt-get install gcc g++ libgomp1", "Compilateurs C++ requis par LightGBM",    BLUE),
    ("COPY requirements.txt → pip install", "Dépendances runtime uniquement",            TEAL),
    ("COPY src/ settings/ models/",     "Code source + modèle sérialisé .dill",         GREEN),
    ("EXPOSE 8000",                     "Port de l'API",                                PURPLE),
    ("HEALTHCHECK urllib /health",      "Surveillance santé du container",              ORANGE),
    ("CMD uvicorn --workers 2",         "2 workers → parallélisation des requêtes",     RED),
]
for i, (cmd, desc, col) in enumerate(layers):
    top = Inches(1.75) + i * Inches(0.72)
    R(s, Inches(0.35), top, Inches(5.8), Inches(0.64), fill=LGREY)
    R(s, Inches(0.35), top, Inches(0.065), Inches(0.64), fill=col)
    T(s, cmd,  Inches(0.55), top+Inches(0.06), Inches(2.8), Inches(0.38), sz=11, bold=True, col=DARK)
    T(s, desc, Inches(3.45), top+Inches(0.06), Inches(2.55), Inches(0.38), sz=11, col=DGREY, italic=True)

T(s, "Commandes", Inches(6.5), Inches(1.25), Inches(6.5), Inches(0.42),
  sz=15, bold=True, col=NAVY)
cmds = [
    ("Build",        "docker build -t laplace-immo-api .",       NAVY),
    ("Run",          "docker run -p 8000:8000 laplace-immo-api", BLUE),
    ("Compose",      "docker compose up",                        TEAL),
    ("Test santé",   "curl http://localhost:8000/health",        GREEN),
    ("Swagger UI",   "http://localhost:8000/docs",               ORANGE),
]
for i, (label, cmd, col) in enumerate(cmds):
    top = Inches(1.75) + i * Inches(0.92)
    R(s, Inches(6.5), top, Inches(6.5), Inches(0.82), fill=LGREY)
    R(s, Inches(6.5), top, Inches(0.065), Inches(0.82), fill=col)
    T(s, label, Inches(6.7), top+Inches(0.08), Inches(1.4), Inches(0.3), sz=11, bold=True, col=DGREY)
    T(s, cmd,   Inches(6.7), top+Inches(0.42), Inches(6.1), Inches(0.3), sz=12, bold=True, col=DARK)

R(s, Inches(6.5), Inches(6.45), Inches(6.5), Inches(0.65), fill=LGREY)
R(s, Inches(6.5), Inches(6.45), Inches(0.065), Inches(0.65), fill=PURPLE)
T(s, "volumes: ./models:/app/models:ro\n→ mise à jour du modèle SANS rebuild image Docker",
  Inches(6.7), Inches(6.52), Inches(6.1), Inches(0.52), sz=11, col=DARK)

# ═══════════════════════════════════════════════════════════════════════════
# 22 — API déployée — Swagger UI en production (Railway)
# ═══════════════════════════════════════════════════════════════════════════
s = S()
HDR(s, "API en Production — Swagger UI (Railway)", "https://clever-harmony-production-ced2.up.railway.app/docs", RED)

swagger_path = Path(__file__).parent.parent / "reports" / "swagger_docs.png"
if swagger_path.exists():
    from pptx.util import Inches as _I
    pic = s.shapes.add_picture(str(swagger_path), Inches(0.35), Inches(1.2), Inches(8.8), Inches(5.6))

T(s, "URL publique", Inches(9.4), Inches(1.2), Inches(3.6), Inches(0.38), sz=13, bold=True, col=NAVY)
T(s, "clever-harmony-production-ced2\n.up.railway.app",
  Inches(9.4), Inches(1.6), Inches(3.6), Inches(0.7), sz=10, col=DGREY, italic=True)

infos = [
    (GREEN,  "/health",         "Statut + modèle"),
    (BLUE,   "/model-info",     "Métriques RMSE/R²"),
    (ORANGE, "/predict",        "Prédiction unitaire"),
    (PURPLE, "/predict/batch",  "Batch 100 maisons"),
]
for i, (col, route, desc) in enumerate(infos):
    top = Inches(2.5) + i * Inches(0.85)
    R(s, Inches(9.4), top, Inches(3.6), Inches(0.76), fill=LGREY)
    R(s, Inches(9.4), top, Inches(0.065), Inches(0.76), fill=col)
    T(s, route, Inches(9.6), top+Inches(0.07), Inches(3.2), Inches(0.3), sz=12, bold=True, col=DARK)
    T(s, desc,  Inches(9.6), top+Inches(0.4), Inches(3.2), Inches(0.28), sz=10, col=DGREY, italic=True)

R(s, Inches(9.4), Inches(6.0), Inches(3.6), Inches(0.72), fill=NAVY)
T(s, "Déployé via CI/CD\nGitHub Actions → Railway",
  Inches(9.5), Inches(6.08), Inches(3.4), Inches(0.55), sz=11, bold=True, col=WHITE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SECTION 8 — MLOPS
# ═══════════════════════════════════════════════════════════════════════════
DIVIDER(NAVY, "MLOps & Monitoring", "⚙")

# 22 — MLFlow & DVC
s = S()
HDR(s, "MLFlow & DVC — Traçabilité & Reproductibilité", "528 runs · 4 expériences · pipeline versionné", NAVY)

T(s, "MLFlow — Tracking complet", Inches(0.35), Inches(1.25), Inches(6.0), Inches(0.42),
  sz=15, bold=True, col=NAVY)
BULS(s, [
    "528 runs enregistrés dans 4 expériences",
    "Par run : hyperparamètres + métriques + artefacts",
    "Nested runs Optuna : trial = run enfant visible",
    "Comparaison visuelle de tous les modèles",
    "Interface : mlflow ui --port 5001",
    "  → http://127.0.0.1:5001",
], Inches(0.35), Inches(1.75), Inches(5.8), sz=13, gap=Inches(0.48))

for i, (exp, runs) in enumerate([
    ("Benchmark initial", "54 runs"),
    ("Optuna LightGBM", "60 trials"),
    ("Optuna XGBoost", "60 trials"),
    ("Pipeline DVC", "n runs"),
]):
    top = Inches(4.2) + i * Inches(0.68)
    R(s, Inches(0.35), top, Inches(5.8), Inches(0.6), fill=LGREY)
    R(s, Inches(0.35), top, Inches(0.065), Inches(0.6), fill=[NAVY,ORANGE,BLUE,GREEN][i])
    T(s, exp, Inches(0.55), top+Inches(0.1), Inches(4.0), Inches(0.38), sz=12, col=DARK)
    T(s, runs, Inches(4.7), top+Inches(0.1), Inches(1.3), Inches(0.38),
      sz=12, bold=True, col=[NAVY,ORANGE,BLUE,GREEN][i], align=PP_ALIGN.RIGHT)

T(s, "DVC — Pipeline reproductible", Inches(6.5), Inches(1.25), Inches(6.5), Inches(0.42),
  sz=15, bold=True, col=NAVY)
BULS(s, [
    "Versionnage des données et des modèles",
    "dvc repro : ré-exécute seulement les étapes impactées",
    "dvc metrics diff : compare RMSE entre branches git",
    "Reproductibilité garantie sur n'importe quelle machine",
    "Compatible git : hash des fichiers = empreinte DVC",
], Inches(6.5), Inches(1.75), Inches(6.5), sz=13, gap=Inches(0.5))

for i, (cmd, out) in enumerate([
    ("make_dataset.py", "data/output/\nhouse_prices.parquet"),
    ("train_pipeline.py", "models/*.dill\nmetrics.json"),
]):
    l = Inches(6.5) + i * Inches(3.3)
    R(s, l, Inches(4.3), Inches(2.9), Inches(1.6), fill=LGREY)
    R(s, l, Inches(4.3), Inches(2.9), Inches(0.065), fill=NAVY)
    T(s, cmd, l+Inches(0.18), Inches(4.4), Inches(2.55), Inches(0.38), sz=11, bold=True, col=DARK)
    T(s, out, l+Inches(0.18), Inches(4.82), Inches(2.55), Inches(0.98), sz=11, col=DGREY)
    if i == 0:
        T(s, "→", l+Inches(2.95), Inches(4.7), Inches(0.4), Inches(0.5), sz=22, bold=True, col=ORANGE)

R(s, Inches(6.5), Inches(6.05), Inches(6.5), Inches(0.8), fill=DARK)
T(s, 'metrics.json → { "rmse_test": 26337.89,  "r2_test": 0.891 }',
  Inches(6.7), Inches(6.2), Inches(6.1), Inches(0.5), sz=12, bold=True, col=RGBColor(0x88,0xFF,0xAA))

# 23 — Monitoring
s = S()
HDR(s, "Monitoring — Détection de Drift", "src/monitoring.py · Evidently + KS-test scipy · 8 features surveillées", NAVY)

T(s, "Pourquoi monitorer en production ?", Inches(0.35), Inches(1.25), Inches(12.5), Inches(0.42),
  sz=15, bold=True, col=NAVY)
BULS(s, [
    "Les données en production dérivent : inflation, nouvelles tendances, saisonnalité",
    "Un modèle précis peut se dégrader silencieusement — le RMSE grimpe sans alerte",
    "Détection précoce → réentraînement ciblé avant que la qualité ne chute",
], Inches(0.35), Inches(1.75), Inches(12.5), sz=13, gap=Inches(0.48))

for i, (title, items, col) in enumerate([
    ("Evidently (rapport complet)", [
        "Rapport HTML interactif",
        "DataDriftPreset : drift par feature",
        "DataQualityPreset : qualité des données",
        "TargetDriftPreset : drift de SalePrice",
        "Sauvegarde : reports/monitoring_drift_report.html",
    ], NAVY),
    ("KS-test scipy (léger)", [
        "Test Kolmogorov-Smirnov par feature",
        "p-value < 0.05 → drift détecté",
        "Rapide, sans dépendance lourde",
        "Fallback si Evidently non disponible",
        "Retourne statistiques détaillées par feature",
    ], BLUE),
    ("Données simulées", [
        "30% des features × facteur 1.2–1.5",
        "Simule un drift de production",
        "En prod : données reçues via l'API",
        "2/8 drifts détectés ✓",
        "bsmtfinsf1 · building_age",
    ], GREEN),
]):
    l = Inches(0.35) + i * Inches(4.35)
    R(s, l, Inches(2.85), Inches(4.1), Inches(3.5), fill=LGREY)
    R(s, l, Inches(2.85), Inches(0.065), Inches(4.1), fill=col)
    T(s, title, l+Inches(0.18), Inches(2.95), Inches(3.8), Inches(0.42), sz=13, bold=True, col=DARK)
    for j, item in enumerate(items):
        T(s, f"▸  {item}", l+Inches(0.18), Inches(3.45)+j*Inches(0.52),
          Inches(3.8), Inches(0.45), sz=11, col=DARK)

R(s, Inches(0.35), Inches(6.55), Inches(12.6), Inches(0.8), fill=LGREY)
R(s, Inches(0.35), Inches(6.55), Inches(0.065), Inches(0.8), fill=GREEN)
T(s, "Résultat test :  2/8 features driftées  →  bsmtfinsf1 (moy. 443→638)  ·  building_age (moy. 36→50)",
  Inches(0.55), Inches(6.68), Inches(12.2), Inches(0.52), sz=13, bold=True, col=DARK)

# 24 — Tests unitaires
s = S()
HDR(s, "Tests Unitaires & Intégration", "20 tests · pytest · mock sans appel réseau", NAVY)

T(s, "20 tests répartis en 3 fichiers", Inches(0.35), Inches(1.25), Inches(12.5), Inches(0.42),
  sz=15, bold=True, col=NAVY)

test_files = [
    ("test_make_dataset.py", "5 tests", BLUE, [
        ("test_load_data_returns_dataframe",   "Retourne un pd.DataFrame"),
        ("test_load_data_lowercase_columns",   "Colonnes en minuscules"),
        ("test_load_data_feature_engineering", "building_age, remodel_age, garage_age créées"),
        ("test_load_data_raises_on_empty_name","ValueError si dataset_name = ''"),
        ("test_load_data_raises_on_none_name", "ValueError si dataset_name = None"),
    ]),
    ("test_trainer.py", "6 tests", TEAL, [
        ("test_define_pipeline_returns_pipeline","define_pipeline() → sklearn Pipeline"),
        ("test_fit_splits_data",                "Après fit(), X_train et X_test non nuls"),
        ("test_train_returns_metrics",          "train() retourne RMSE, MAE, R²"),
        ("test_predict_returns_correct_shape",  "predict() retourne la bonne longueur"),
        ("test_train_test_size_respected",      "test_size = 20% respecté"),
        ("test_score_returns_float",            "score() retourne un float"),
    ]),
    ("test_api.py", "9 tests", ORANGE, [
        ("test_root",                           "GET / → 200 avec clé 'docs'"),
        ("test_health",                         "model_loaded = True"),
        ("test_predict_returns_positive_price", "prix prédit > 0"),
        ("test_predict_high_quality_house_costs_more", "qualité 9 > qualité 3 (sanity)"),
        ("test_predict_batch",                  "Batch 2 maisons → count=2"),
        ("test_docs_available",                 "Swagger UI accessible"),
        ("test_openapi_schema",                 "/predict dans le schéma OpenAPI"),
    ]),
]
for i, (fname, count, col, tests) in enumerate(test_files):
    l = Inches(0.35) + i * Inches(4.35)
    R(s, l, Inches(1.75), Inches(4.1), Inches(5.7), fill=LGREY)
    R(s, l, Inches(1.75), Inches(0.065), Inches(4.1), fill=col)
    R(s, l+Inches(2.9), Inches(1.82), Inches(0.9), Inches(0.38), fill=col)
    T(s, fname, l+Inches(0.18), Inches(1.87), Inches(2.65), Inches(0.42), sz=12, bold=True, col=DARK)
    T(s, count, l+Inches(2.9), Inches(1.82), Inches(0.9), Inches(0.38),
      sz=12, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
    for j, (test_name, desc) in enumerate(tests[:6]):
        top = Inches(2.28) + j * Inches(0.56)
        T(s, "✓", l+Inches(0.18), top, Inches(0.28), Inches(0.45), sz=12, col=GREEN, bold=True)
        T(s, test_name, l+Inches(0.5), top+Inches(0.02), Inches(3.4), Inches(0.3), sz=10, bold=True, col=DARK)
        T(s, desc, l+Inches(0.5), top+Inches(0.3), Inches(3.4), Inches(0.22), sz=9, col=DGREY)

R(s, Inches(0.35), Inches(6.7), Inches(12.6), Inches(0.75), fill=LGREY)
R(s, Inches(0.35), Inches(6.7), Inches(12.6), Inches(0.065), fill=GREEN)
T(s, "Tous les tests utilisent unittest.mock.patch pour simuler fetch_openml — aucun appel réseau",
  Inches(0.55), Inches(6.82), Inches(12.2), Inches(0.5), sz=13, col=DARK, italic=True)

# 25 — CI/CD
s = S()
HDR(s, "CI/CD — GitHub Actions", "3 jobs séquentiels · automatique à chaque push sur n'importe quelle branche", NAVY)

T(s, "Pipeline automatique", Inches(0.35), Inches(1.25), Inches(12.5), Inches(0.42),
  sz=15, bold=True, col=NAVY)

jobs_detail = [
    ("Job 1", "Tests Unitaires", "pytest tests/ -v\n--ignore=tests/test_api.py\n\n11 tests (make_dataset + trainer)\n\nSi échec → Jobs 2 & 3 bloqués", BLUE),
    ("Job 2", "Pipeline DVC", "needs: tests\n\npython src/train_pipeline.py\n\nEntraîne le modèle\nUpload metrics.json\ncomme artefact GitHub", GREEN),
    ("Job 3", "Tests API + Docker", "needs: dvc-pipeline\n\npython src/train_pipeline.py\npytest tests/test_api.py\n  (9 tests intégration)\ndocker build -t laplace-immo-api .", ORANGE),
]

for i, (num, title, desc, col) in enumerate(jobs_detail):
    l = Inches(0.35) + i * Inches(4.35)
    R(s, l, Inches(1.75), Inches(4.1), Inches(4.8), fill=LGREY)
    R(s, l, Inches(1.75), Inches(0.065), Inches(4.1), fill=col)
    T(s, num, l+Inches(0.18), Inches(1.85), Inches(0.8), Inches(0.45), sz=14, bold=True, col=col)
    T(s, title, l+Inches(1.05), Inches(1.85), Inches(2.9), Inches(0.45), sz=14, bold=True, col=DARK)
    T(s, desc, l+Inches(0.18), Inches(2.38), Inches(3.75), Inches(4.0), sz=11, col=DARK)
    if i < 2:
        T(s, "↓ si succès", l+Inches(4.15), Inches(3.6), Inches(0.5), Inches(0.8),
          sz=11, col=MGREY, bold=True, align=PP_ALIGN.CENTER)

T(s, "Déclencheur", Inches(0.35), Inches(6.7), Inches(3.0), Inches(0.4),
  sz=13, bold=True, col=NAVY)
T(s, "on: push (branches: **) + pull_request (branches: **) — tous les push déclenchent la CI",
  Inches(0.35), Inches(7.0), Inches(12.5), Inches(0.4), sz=12, col=DARK, italic=True)

# ═══════════════════════════════════════════════════════════════════════════
# 26 — STACK
# ═══════════════════════════════════════════════════════════════════════════
s = S()
HDR(s, "Stack Technique Complète", "")

T(s, "Toutes les technologies utilisées dans ce projet", Inches(0.35), Inches(1.25), Inches(12.5), Inches(0.42),
  sz=15, bold=True, col=NAVY)

stack = [
    ("Machine Learning",   "scikit-learn · XGBoost · LightGBM · CatBoost · Ridge · SVR",          NAVY),
    ("Optimisation",       "Optuna 4.4 (TPE, pruning Median) · 60 trials/modèle",                  ORANGE),
    ("MLOps",              "MLFlow 3.1 (528 runs, 4 expériences) · DVC 3+ (pipeline reproductible)",GREEN),
    ("Feature Engineering","PPS · TargetEncoder · Yeo-Johnson · OrdinalEncoder · Interactions",     TEAL),
    ("API",                "FastAPI · Pydantic · uvicorn[standard] · httpx",                        RED),
    ("Sérialisation",      "dill 0.4 (supporte lambdas + objets Python complexes)",                 PURPLE),
    ("Monitoring",         "Evidently · scipy (KS-test) · category-encoders",                       BLUE),
    ("Containerisation",   "Docker (python:3.11-slim) · docker-compose",                            TEAL),
    ("Tests",              "pytest 8+ · pytest-asyncio · FastAPI TestClient · unittest.mock",       GREEN),
    ("CI/CD",              "GitHub Actions (4 jobs) · Railway déploiement automatique",            NAVY),
    ("Utilitaires",        "loguru · pendulum · pyarrow · numpy 1.26 · pandas 2.3",                 DGREY),
]
for i, (cat, tools, col) in enumerate(stack):
    top = Inches(1.75) + i * Inches(0.52)
    R(s, Inches(0.35), top, Inches(12.6), Inches(0.46), fill=LGREY if i%2==0 else WHITE)
    R(s, Inches(0.35), top, Inches(0.065), Inches(0.46), fill=col)
    T(s, cat, Inches(0.55), top+Inches(0.07), Inches(2.3), Inches(0.35), sz=12, bold=True, col=DARK)
    T(s, tools, Inches(2.95), top+Inches(0.07), Inches(10.0), Inches(0.35), sz=11, col=DGREY)

# ═══════════════════════════════════════════════════════════════════════════
# 27 — CONCLUSION
# ═══════════════════════════════════════════════════════════════════════════
s = S()
R(s, 0, 0, W, Inches(1.05), fill=NAVY)
R(s, 0, Inches(1.05), W, Inches(0.055), fill=ORANGE)
T(s, "Bilan & Conclusion", Inches(0.45), Inches(0.15), W-Inches(0.9), Inches(0.75),
  sz=30, bold=True, col=WHITE)

KPI(s, Inches(0.35), Inches(1.25), Inches(2.4), Inches(1.6), "26 338 $", "RMSE Test",   "erreur prédiction",    NAVY)
KPI(s, Inches(2.9),  Inches(1.25), Inches(2.4), Inches(1.6), "0.891",    "R² Test",     "variance expliquée",   GREEN)
KPI(s, Inches(5.45), Inches(1.25), Inches(2.4), Inches(1.6), "17 058 $", "MAE Test",    "écart absolu moyen",   ORANGE)
KPI(s, Inches(7.99), Inches(1.25), Inches(2.4), Inches(1.6), "20",       "Tests",       "unitaires + intégr.",  TEAL)
KPI(s, Inches(10.55),Inches(1.25), Inches(2.44),Inches(1.6), "528",      "Runs MLFlow", "4 expériences",        PURPLE)

T(s, "Ce que nous avons accompli", Inches(0.35), Inches(3.1), Inches(6.0), Inches(0.42),
  sz=14, bold=True, col=NAVY)
BULS(s, [
    "EDA rigoureuse → découverte log-transform (+54 000 $ sur SVR)",
    "Benchmark exhaustif : 20+ modèles, 54 runs MLFlow",
    "Optimisation Optuna : 60 trials, pruning bayésien",
    "6 axes de preprocessing (+1 237 $ de gain total)",
    "API REST déployée sur Railway (URL publique)",
    "Monitoring de drift avec Evidently + KS-test",
    "20 tests automatisés + CI/CD 4 jobs GitHub Actions",
], Inches(0.35), Inches(3.6), Inches(6.0), sz=12, gap=Inches(0.46))

T(s, "Perspectives d'amélioration", Inches(6.5), Inches(3.1), Inches(6.5), Inches(0.42),
  sz=14, bold=True, col=NAVY)
BULS(s, [
    "SHAP values : explicabilité des prédictions par maison",
    "Réentraînement automatique si drift détecté (Airflow/Prefect)",
    "Test A/B entre modèle Ridge et LightGBM en production",
    "Interface web Streamlit pour les agents immobiliers",
    "Étendre à d'autres marchés immobiliers (Paris, NY…)",
], Inches(6.5), Inches(3.6), Inches(6.5), sz=12, gap=Inches(0.46))

R(s, Inches(0.35), Inches(6.8), Inches(12.6), Inches(0.62), fill=NAVY)
T(s, "Lykarim  &  AlimatouSadiyah  ·  Promotion DIC3  ·  2026  ·  GitHub : Lykarim/house-price-ames",
  Inches(0.35), Inches(6.9), Inches(12.6), Inches(0.42),
  sz=14, bold=True, col=WHITE, align=PP_ALIGN.CENTER)

# ── Sauvegarde ────────────────────────────────────────────────────────────
out = Path(__file__).parent / "laplace_immo_presentation.pptx"
prs.save(str(out))
n = len(prs.slides)
print(f"✓ {n} slides → {out}")
